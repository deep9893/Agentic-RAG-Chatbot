# utils.py
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from docx import Document

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    return " ".join(page.get_text() for page in doc)

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_from_csv(path):
    df = pd.read_csv(path)
    return df.to_string()

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()
